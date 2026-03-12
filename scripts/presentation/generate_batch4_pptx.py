#!/usr/bin/env python3
"""
Generate Batch 4 slides (14-17) using python-pptx with exact Design System.
Hackathon PoD Academy — Claro + Oracle
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# =============================================================================
# DESIGN SYSTEM CONSTANTS
# =============================================================================
SLIDE_WIDTH = Inches(13.333)   # 1920px @ 144dpi → 16:9
SLIDE_HEIGHT = Inches(7.5)     # 1080px

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_CARD = RGBColor(0xF8, 0xF9, 0xFA)
CYAN = RGBColor(0x00, 0xF0, 0xFF)
GREEN = RGBColor(0x00, 0xFF, 0x88)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)
RED = RGBColor(0xFF, 0x33, 0x66)
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
BODY_COLOR = RGBColor(0x4A, 0x4A, 0x6A)
LABEL_COLOR = RGBColor(0x8A, 0x8A, 0xAA)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
BRONZE_COLOR = RGBColor(0xCD, 0x7F, 0x32)
CYAN_30 = RGBColor(0xB3, 0xFB, 0xFF)  # cyan at ~30% on white
TEAL_DARK = RGBColor(0x00, 0xC4, 0xD4)

MARGIN = Inches(0.6)  # ~80px
HEADER_H = Inches(0.55)
FOOTER_H = Inches(0.35)
CONTENT_TOP = HEADER_H + Inches(0.15)
CONTENT_BOTTOM = SLIDE_HEIGHT - FOOTER_H - Inches(0.1)
CONTENT_W = SLIDE_WIDTH - 2 * MARGIN


def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================
def add_circuit_decoration(slide, prs):
    """Add subtle circuit-like decorative lines in corners."""
    line_color = CYAN
    line_w = Pt(0.75)
    corner_len = Inches(1.2)

    corners = [
        # top-left
        (Inches(0.2), Inches(0.15), corner_len, Pt(0), True),
        (Inches(0.2), Inches(0.15), Pt(0), corner_len, False),
        # top-right
        (SLIDE_WIDTH - Inches(0.2) - corner_len, Inches(0.15), corner_len, Pt(0), True),
        (SLIDE_WIDTH - Inches(0.2), Inches(0.15), Pt(0), corner_len, False),
        # bottom-left
        (Inches(0.2), SLIDE_HEIGHT - Inches(0.15), corner_len, Pt(0), True),
        (Inches(0.2), SLIDE_HEIGHT - Inches(0.15) - corner_len, Pt(0), corner_len, False),
        # bottom-right
        (SLIDE_WIDTH - Inches(0.2) - corner_len, SLIDE_HEIGHT - Inches(0.15), corner_len, Pt(0), True),
        (SLIDE_WIDTH - Inches(0.2), SLIDE_HEIGHT - Inches(0.15) - corner_len, Pt(0), corner_len, False),
    ]

    for left, top, w, h, is_horizontal in corners:
        if is_horizontal:
            connector = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, w, Pt(1.5)
            )
        else:
            connector = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, Pt(1.5), h
            )
        connector.fill.solid()
        connector.fill.fore_color.rgb = CYAN
        connector.fill.fore_color.brightness = 0.7  # 30% opacity effect
        connector.line.fill.background()


def add_header(slide, prs):
    """Add header bar with Claro + Oracle text and cyan bottom border."""
    # Header background
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_WIDTH, HEADER_H
    )
    header.fill.solid()
    header.fill.fore_color.rgb = WHITE
    header.line.fill.background()

    # Claro text (left)
    tf_claro = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.08), Inches(2), Inches(0.4)
    )
    p = tf_claro.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Claro"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.name = "Inter"

    # Oracle text (right)
    tf_oracle = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(2.5), Inches(0.08), Inches(2), Inches(0.4)
    )
    p = tf_oracle.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "ORACLE"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xC7, 0x41, 0x34)  # Oracle red
    run.font.name = "Inter"

    # Cyan bottom border
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), HEADER_H - Pt(3),
        SLIDE_WIDTH, Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()


def add_footer(slide, slide_num, total=17):
    """Add footer with text and slide number."""
    # Left text
    tf_left = slide.shapes.add_textbox(
        MARGIN, SLIDE_HEIGHT - FOOTER_H, Inches(4), Inches(0.3)
    )
    p = tf_left.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Hackathon PoD Academy — Março 2026"
    run.font.size = Pt(10)
    run.font.color.rgb = LABEL_COLOR
    run.font.name = "Inter"

    # Right text
    tf_right = slide.shapes.add_textbox(
        SLIDE_WIDTH - MARGIN - Inches(2), SLIDE_HEIGHT - FOOTER_H,
        Inches(2), Inches(0.3)
    )
    p = tf_right.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Slide {slide_num}/{total}"
    run.font.size = Pt(10)
    run.font.color.rgb = LABEL_COLOR
    run.font.name = "Inter"


def add_slide_title(slide, title, subtitle=None, top=None):
    """Add slide title and optional subtitle."""
    if top is None:
        top = CONTENT_TOP + Inches(0.1)

    tf = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, Inches(0.5))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.name = "Inter"

    sub_top = top + Inches(0.5)
    if subtitle:
        tf2 = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, Inches(0.35))
        tf2.text_frame.word_wrap = True
        p2 = tf2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(16)
        run2.font.color.rgb = BODY_COLOR
        run2.font.name = "Inter"
        return sub_top + Inches(0.4)
    return sub_top


def add_card(slide, left, top, width, height, fill_color=BG_CARD,
             border_color=None, border_width=Pt(1), glow=False):
    """Add a rounded-corner card shape."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = border_width
    else:
        card.line.color.rgb = CYAN
        card.line.width = Pt(1)
        card.line.color.brightness = 0.7
    # Adjust corner rounding
    card.adjustments[0] = 0.05  # ~12px radius effect
    return card


def add_text_in_card(slide, left, top, width, height, lines,
                     alignment=PP_ALIGN.LEFT):
    """Add multi-line text inside a card area.
    lines: list of (text, font_size_pt, color, bold, italic)
    """
    tf = slide.shapes.add_textbox(left, top, width, height)
    tf.text_frame.word_wrap = True

    for i, (text, size, color, bold, italic) in enumerate(lines):
        if i == 0:
            p = tf.text_frame.paragraphs[0]
        else:
            p = tf.text_frame.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(2)
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = "Inter"
    return tf


def add_badge(slide, left, top, text, bg_color, text_color=WHITE, font_size=9):
    """Add a pill-shaped badge."""
    w = Inches(max(0.7, len(text) * 0.08 + 0.3))
    h = Inches(0.22)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_color
    badge.line.fill.background()
    badge.adjustments[0] = 0.5  # full pill shape

    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Inter"
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    return badge


def add_arrow(slide, left, top, width, direction='right'):
    """Add a cyan arrow."""
    if direction == 'right':
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.25)
        )
    elif direction == 'down':
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.25), width
        )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = CYAN
    arrow.line.fill.background()
    return arrow


def add_connecting_line(slide, x1, y1, x2, y2, color=CYAN, width=Pt(1.5)):
    """Add a simple connecting line between two points."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = width


# =============================================================================
# SLIDE 14 — SQUAD OCI
# =============================================================================
def create_slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_header(slide, prs)
    add_footer(slide, 14)
    add_circuit_decoration(slide, prs)

    content_top = add_slide_title(
        slide,
        "Squad OCI — 6 Agentes IA Especializados",
        "Orquestração multi-agente com DNA de 9 especialistas reais"
    )

    # --- ATLAS Card (Orchestrator) - centered top ---
    atlas_w = Inches(5.5)
    atlas_h = Inches(1.65)
    atlas_left = (SLIDE_WIDTH - atlas_w) / 2
    atlas_top = content_top + Inches(0.1)

    add_card(slide, atlas_left, atlas_top, atlas_w, atlas_h,
             fill_color=WHITE, border_color=CYAN, border_width=Pt(2.5))

    add_text_in_card(slide, atlas_left + Inches(0.3), atlas_top + Inches(0.1),
                     atlas_w - Inches(0.6), atlas_h - Inches(0.2), [
        ("☁️  ATLAS", 24, TITLE_COLOR, True, False),
        ("Orquestrador — OCI Platform Chief", 12, BODY_COLOR, False, False),
        ("Elite Mind: Rohit Rahi — VP Oracle University, OCI Launch Team (2016)", 10, LABEL_COLOR, False, False),
        ("Delega, roteia e valida. NUNCA executa trabalho especialista.", 10, BODY_COLOR, False, False),
        ('"Orquestrar, delegar, quality gate, checkpoint"', 9, LABEL_COLOR, False, True),
    ], alignment=PP_ALIGN.CENTER)

    add_badge(slide, (SLIDE_WIDTH - Inches(1.2)) / 2,
              atlas_top + Inches(0.55), "ORCHESTRATOR", CYAN)

    # --- 5 Specialist Cards ---
    agents = [
        {
            "icon": "🏗️", "name": "STRATOS", "role": "Infraestrutura — IaC Architect",
            "minds": ["André Corrêa Neto — CIS OCI Landing Zones",
                      "Ali Mukadam — Terraform OCI Modules",
                      "Yevgeniy Brikman — Terraform Up & Running"],
            "tier": "TIER 0", "tier_color": RED,
            "quote": '"Infraestrutura como código"'
        },
        {
            "icon": "🔄", "name": "FLUX", "role": "Data Pipeline — Lakehouse & PySpark",
            "minds": ["Matei Zaharia — Criador Apache Spark",
                      "Holden Karau — High Performance Spark"],
            "tier": "TIER 1", "tier_color": CYAN,
            "quote": '"Data flows through the medallion"'
        },
        {
            "icon": "🧠", "name": "NEURON", "role": "ML Systems — Credit Risk Specialist",
            "minds": ["Chip Huyen — Designing ML Systems (Stanford)",
                      "Goku Mohandas — Made With ML (30K+ stars)"],
            "tier": "TIER 1", "tier_color": CYAN,
            "quote": '"Modelos em produção na OCI"'
        },
        {
            "icon": "💰", "name": "NIMBUS", "role": "FinOps — Cloud Operations",
            "minds": ["J.R. Storment — FinOps Foundation",
                      "Rohit Rahi — OCI Architecture"],
            "tier": "TIER 1", "tier_color": CYAN,
            "quote": '"Every dollar has a purpose"'
        },
        {
            "icon": "🛡️", "name": "SENTINEL", "role": "Segurança — CIS Compliance",
            "minds": ["André Corrêa Neto — CIS OCI Benchmarks"],
            "tier": "TIER 2", "tier_color": YELLOW, "tier_text_color": TITLE_COLOR,
            "quote": '"Segurança é inegociável"'
        },
    ]

    card_w = Inches(2.2)
    card_h = Inches(2.3)
    gap = Inches(0.2)
    total_w = 5 * card_w + 4 * gap
    start_left = (SLIDE_WIDTH - total_w) / 2
    cards_top = atlas_top + atlas_h + Inches(0.55)

    # Connection lines from Atlas to each card
    atlas_bottom_y = atlas_top + atlas_h
    for i in range(5):
        card_left = start_left + i * (card_w + gap)
        card_center_x = card_left + card_w / 2
        # vertical line
        add_connecting_line(
            slide,
            card_center_x, atlas_bottom_y,
            card_center_x, cards_top - Inches(0.02)
        )

    # Horizontal line connecting all verticals
    first_center = start_left + card_w / 2
    last_center = start_left + 4 * (card_w + gap) + card_w / 2
    mid_y = atlas_bottom_y + Inches(0.2)
    add_connecting_line(slide, first_center, mid_y, last_center, mid_y)

    # Vertical from Atlas center down to horizontal
    atlas_center_x = SLIDE_WIDTH / 2
    add_connecting_line(slide, atlas_center_x, atlas_bottom_y, atlas_center_x, mid_y)

    # Vertical from horizontal down to each card
    for i in range(5):
        card_left = start_left + i * (card_w + gap)
        card_center_x = card_left + card_w / 2
        add_connecting_line(slide, card_center_x, mid_y, card_center_x, cards_top)

    for i, agent in enumerate(agents):
        card_left = start_left + i * (card_w + gap)

        add_card(slide, card_left, cards_top, card_w, card_h,
                 fill_color=BG_CARD, border_color=CYAN)

        # Build text lines
        lines = [
            (f"{agent['icon']}  {agent['name']}", 16, TITLE_COLOR, True, False),
            (agent['role'], 10, BODY_COLOR, False, False),
            ("", 6, BODY_COLOR, False, False),  # spacer
            ("Elite Minds:", 9, LABEL_COLOR, True, False),
        ]
        for mind in agent['minds']:
            lines.append((mind, 8, LABEL_COLOR, False, False))
        lines.append(("", 4, BODY_COLOR, False, False))
        lines.append((agent['quote'], 8, LABEL_COLOR, False, True))

        add_text_in_card(slide, card_left + Inches(0.12), cards_top + Inches(0.1),
                         card_w - Inches(0.24), card_h - Inches(0.5), lines,
                         alignment=PP_ALIGN.CENTER)

        # Tier badge
        tier_text_color = agent.get('tier_text_color', WHITE)
        add_badge(slide, card_left + (card_w - Inches(0.7)) / 2,
                  cards_top + Inches(0.55),
                  agent['tier'], agent['tier_color'], tier_text_color, 8)

    # --- Bottom insight bar ---
    bar_top = cards_top + card_h + Inches(0.15)
    bar_h = Inches(0.35)
    add_card(slide, MARGIN, bar_top, CONTENT_W, bar_h,
             fill_color=BG_CARD, border_color=CYAN)
    add_text_in_card(slide, MARGIN + Inches(0.2), bar_top + Inches(0.03),
                     CONTENT_W - Inches(0.4), bar_h, [
        ("Quality Mode A+  |  ~150 veto conditions  |  9 elite minds clonados  |  Modelo: All Opus — máximo raciocínio", 12, BODY_COLOR, False, False),
    ], alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 15 — MIND CLONING
# =============================================================================
def create_slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, prs)
    add_footer(slide, 15)
    add_circuit_decoration(slide, prs)

    content_top = add_slide_title(
        slide,
        "Mind Cloning — DNA de Especialistas Reais",
        "Pipeline de 3 fases para clonar expertise em agentes IA"
    )

    # 3 columns
    col_w = Inches(3.6)
    col_h = Inches(3.3)
    gap = Inches(0.35)
    arrow_w = Inches(0.45)
    total_w = 3 * col_w + 2 * (gap + arrow_w + gap)
    start_left = (SLIDE_WIDTH - total_w) / 2
    col_top = content_top + Inches(0.15)

    columns = [
        {
            "num": "1", "title": "RESEARCH", "subtitle": "Curadoria de Fontes",
            "icon": "🔍",
        },
        {
            "num": "2", "title": "DNA EXTRACTION", "subtitle": "Voice + Thinking DNA",
            "icon": "🧬",
        },
        {
            "num": "3", "title": "ARTIFACT", "subtitle": "Agente Operacional",
            "icon": "⚙️",
        },
    ]

    for i, col in enumerate(columns):
        col_left = start_left + i * (col_w + gap + arrow_w + gap)

        # Card
        add_card(slide, col_left, col_top, col_w, col_h,
                 fill_color=WHITE, border_color=CYAN, border_width=Pt(1.5))

        # Title
        add_text_in_card(slide, col_left + Inches(0.15), col_top + Inches(0.1),
                         col_w - Inches(0.3), Inches(0.6), [
            (f"{col['icon']}  {col['num']}. {col['title']}", 18, TITLE_COLOR, True, False),
            (col['subtitle'], 12, BODY_COLOR, False, False),
        ])

        # Arrow between columns
        if i < 2:
            arrow_left = col_left + col_w + gap
            arrow_top = col_top + col_h / 2 - Inches(0.12)
            add_arrow(slide, arrow_left, arrow_top, arrow_w)
            # Label above arrow
            label_text = "Extrai" if i == 0 else "Constrói"
            tf = slide.shapes.add_textbox(
                arrow_left - Inches(0.1), arrow_top - Inches(0.22),
                arrow_w + Inches(0.2), Inches(0.2)
            )
            p = tf.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label_text
            run.font.size = Pt(9)
            run.font.color.rgb = CYAN
            run.font.bold = True
            run.font.name = "Inter"

    # --- Column 1 content: Source quality tiers ---
    c1_left = start_left
    c1_content_top = col_top + Inches(0.75)
    c1_pad = Inches(0.15)

    content_lines_c1 = [
        ("• Pesquisa de especialistas reais por domínio", 10, BODY_COLOR, False, False),
        ("• Curadoria de fontes primárias (livros, papers, talks)", 10, BODY_COLOR, False, False),
        ("• Classificação por qualidade:", 10, BODY_COLOR, False, False),
    ]
    add_text_in_card(slide, c1_left + c1_pad, c1_content_top,
                     col_w - 2 * c1_pad, Inches(0.7), content_lines_c1)

    # Mini-cards for tiers
    tier_top = c1_content_top + Inches(0.7)
    tiers = [
        ("OURO — Fonte ≥60% coverage", "Chip Huyen → Designing ML Systems", GREEN),
        ("BRONZE — Fonte parcial", "Blog posts, conference talks", CYAN),
        ("ELIMINATE — Não confiável", "Resumos de terceiros, conteúdo desatualizado", RED),
    ]
    for j, (tier_label, tier_ex, tier_col) in enumerate(tiers):
        t_top = tier_top + j * Inches(0.55)
        mini = add_card(slide, c1_left + Inches(0.2), t_top,
                        col_w - Inches(0.4), Inches(0.48),
                        fill_color=WHITE, border_color=tier_col, border_width=Pt(2))
        add_text_in_card(slide, c1_left + Inches(0.3), t_top + Inches(0.03),
                         col_w - Inches(0.6), Inches(0.42), [
            (tier_label, 9, TITLE_COLOR, True, False),
            (tier_ex, 8, LABEL_COLOR, False, True),
        ])

    # --- Column 2 content: Voice DNA + Thinking DNA ---
    c2_left = start_left + (col_w + gap + arrow_w + gap)
    c2_content_top = col_top + Inches(0.75)
    sub_w = (col_w - Inches(0.4)) / 2

    # Voice DNA sub-card
    add_card(slide, c2_left + Inches(0.12), c2_content_top,
             sub_w, Inches(2.2), fill_color=BG_CARD, border_color=CYAN)
    add_text_in_card(slide, c2_left + Inches(0.2), c2_content_top + Inches(0.05),
                     sub_w - Inches(0.16), Inches(2.1), [
        ("Voice DNA", 12, TITLE_COLOR, True, False),
        ("Signature phrases", 9, BODY_COLOR, False, False),
        ("Power words", 9, BODY_COLOR, False, False),
        ("Tone calibration", 9, BODY_COLOR, False, False),
        ("Prohibited phrases", 9, BODY_COLOR, False, False),
        ("", 6, BODY_COLOR, False, False),
        ('Zaharia: "The medallion pattern exists because data matures in stages"', 7, LABEL_COLOR, False, True),
        ("", 4, BODY_COLOR, False, False),
        ('Storment: "Every dollar has a purpose"', 7, LABEL_COLOR, False, True),
    ])

    # Thinking DNA sub-card
    add_card(slide, c2_left + Inches(0.12) + sub_w + Inches(0.12), c2_content_top,
             sub_w, Inches(2.2), fill_color=BG_CARD, border_color=CYAN)
    add_text_in_card(slide, c2_left + Inches(0.24) + sub_w + Inches(0.08), c2_content_top + Inches(0.05),
                     sub_w - Inches(0.16), Inches(2.1), [
        ("Thinking DNA", 12, TITLE_COLOR, True, False),
        ("Decision frameworks", 9, BODY_COLOR, False, False),
        ("Heuristics & SOPs", 9, BODY_COLOR, False, False),
        ("Anti-patterns catalog", 9, BODY_COLOR, False, False),
        ("Veto conditions", 9, BODY_COLOR, False, False),
        ("", 6, BODY_COLOR, False, False),
        ('Huyen: "A model is only as good as the system that serves it"', 7, LABEL_COLOR, False, True),
        ("", 4, BODY_COLOR, False, False),
        ('Brikman: "Small, composable, testable, versioned modules"', 7, LABEL_COLOR, False, True),
    ])

    # --- Column 3 content: Artifact checklist ---
    c3_left = start_left + 2 * (col_w + gap + arrow_w + gap)
    c3_content_top = col_top + Inches(0.75)

    checklist = [
        "✅  Agent persona file (.md)",
        "✅  Core principles (CRITICAL/HIGH)",
        "✅  Commands & routing matrix",
        "✅  Quality gates por fase",
        "✅  Veto conditions (~150 total)",
        "✅  Anti-patterns catalog",
        "✅  Psychometric calibration",
    ]
    lines_c3 = [(item, 10, BODY_COLOR, False, False) for item in checklist]
    add_text_in_card(slide, c3_left + Inches(0.2), c3_content_top,
                     col_w - Inches(0.4), Inches(1.6), lines_c3)

    # Result mini-card
    result_top = c3_content_top + Inches(1.7)
    add_card(slide, c3_left + Inches(0.15), result_top,
             col_w - Inches(0.3), Inches(0.6),
             fill_color=WHITE, border_color=GREEN, border_width=Pt(2))
    add_text_in_card(slide, c3_left + Inches(0.25), result_top + Inches(0.05),
                     col_w - Inches(0.5), Inches(0.5), [
        ("Resultado: Agente que PENSA como o especialista, NÃO apenas repete frases", 9, TITLE_COLOR, True, False),
    ], alignment=PP_ALIGN.CENTER)

    # --- Bottom metrics bar ---
    bar_top = col_top + col_h + Inches(0.25)
    bar_h = Inches(0.65)
    kpis = [
        ("9 Elite Minds", "Especialistas clonados"),
        ("6 Agentes", "Operacionais no squad"),
        ("~150 Veto Conditions", "Quality Mode A+"),
        ("12 Anti-patterns", "Catalogados e bloqueados"),
    ]
    kpi_w = CONTENT_W / 4
    for i, (val, label) in enumerate(kpis):
        kpi_left = MARGIN + i * kpi_w
        add_card(slide, kpi_left + Inches(0.05), bar_top,
                 kpi_w - Inches(0.1), bar_h, fill_color=BG_CARD)
        add_text_in_card(slide, kpi_left + Inches(0.1), bar_top + Inches(0.05),
                         kpi_w - Inches(0.2), bar_h - Inches(0.1), [
            (val, 18, TITLE_COLOR, True, False),
            (label, 9, LABEL_COLOR, False, False),
        ], alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 16 — PRODUTIVIDADE IA
# =============================================================================
def create_slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, prs)
    add_footer(slide, 16)
    add_circuit_decoration(slide, prs)

    content_top = add_slide_title(
        slide,
        "Produtividade IA — Resultados Concretos",
        "O que a orquestração multi-agente entregou no hackathon"
    )

    # --- Timeline bar (top) ---
    timeline_top = content_top + Inches(0.1)
    phases = [
        ("INFRA", "Terraform\n7 módulos", DARK_BG),
        ("DATA", "Pipeline B+S+G\n5h45", DARK_BG),
        ("ML", "Training 5 modelos\n58 min", DARK_BG),
        ("CONSUME", "ADW+APEX\nAlways Free", DARK_BG),
        ("TOTAL", "~6h43 pipeline\ncompleto", CYAN),
    ]
    phase_w = CONTENT_W / 5 - Inches(0.08)
    phase_h = Inches(0.75)
    phase_gap = Inches(0.1)

    for i, (name, desc, bg) in enumerate(phases):
        p_left = MARGIN + i * (phase_w + phase_gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, p_left, timeline_top, phase_w, phase_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg
        card.line.fill.background()
        card.adjustments[0] = 0.08

        text_color = WHITE if bg == DARK_BG else TITLE_COLOR
        tf = slide.shapes.add_textbox(
            p_left + Inches(0.08), timeline_top + Inches(0.05),
            phase_w - Inches(0.16), phase_h - Inches(0.1)
        )
        tf.text_frame.word_wrap = True
        p = tf.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = text_color
        run.font.name = "Inter"

        p2 = tf.text_frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(9)
        run2.font.color.rgb = text_color if bg == CYAN else RGBColor(0xCC, 0xCC, 0xCC)
        run2.font.name = "Inter"

        # Arrow between phases
        if i < 4:
            arrow_left = p_left + phase_w
            arrow_top_pos = timeline_top + phase_h / 2 - Inches(0.08)
            a = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top_pos,
                phase_gap, Inches(0.16)
            )
            a.fill.solid()
            a.fill.fore_color.rgb = CYAN
            a.line.fill.background()

    # --- 6 KPI Cards (2 rows x 3 cols) ---
    kpi_top = timeline_top + phase_h + Inches(0.3)
    kpis = [
        ("6", "Scripts Python", "Produção-ready, testados", TITLE_COLOR),
        ("13", "Tasks Airflow", "DAG orquestrado end-to-end", TITLE_COLOR),
        ("7", "Módulos Terraform", "IaC completo, destroy < 5min", TITLE_COLOR),
        ("5", "Modelos Treinados", "Todos QG-05 PASSED", TITLE_COLOR),
        ("3.9M", "Registros Scorados", "Batch scoring completo", TITLE_COLOR),
        ("0", "Erros de Produção", "~150 veto conditions", GREEN),
    ]

    kpi_cols = 6
    kpi_w = CONTENT_W / kpi_cols - Inches(0.08)
    kpi_h = Inches(1.15)

    for i, (val, label, sub, val_color) in enumerate(kpis):
        col = i % kpi_cols
        k_left = MARGIN + col * (kpi_w + Inches(0.08))
        k_top = kpi_top

        add_card(slide, k_left, k_top, kpi_w, kpi_h,
                 fill_color=WHITE, border_color=CYAN)

        add_text_in_card(slide, k_left + Inches(0.08), k_top + Inches(0.08),
                         kpi_w - Inches(0.16), kpi_h - Inches(0.16), [
            (val, 36, val_color, True, False),
            (label, 11, TITLE_COLOR, True, False),
            (sub, 8, LABEL_COLOR, False, False),
        ], alignment=PP_ALIGN.CENTER)

    # --- Tools section ---
    tools_top = kpi_top + kpi_h + Inches(0.3)
    tools_title = slide.shapes.add_textbox(
        MARGIN, tools_top, CONTENT_W, Inches(0.3)
    )
    p = tools_title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Ferramentas dos Agentes"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.name = "Inter"

    tool_rows = [
        ("Infra + Ops", ["Terraform", "OCI CLI", "Steampipe", "Checkov", "Resource Scheduler"]),
        ("Data + ML", ["PySpark 3.5", "Delta/Parquet", "MLflow", "Optuna", "SHAP", "Pandera"]),
        ("Plataforma", ["OCI Data Flow", "Object Storage", "ADW Serverless", "Data Catalog", "OCI Vault", "APEX"]),
    ]

    row_top = tools_top + Inches(0.35)
    for row_i, (row_label, tools) in enumerate(tool_rows):
        r_top = row_top + row_i * Inches(0.35)

        # Row label
        tf = slide.shapes.add_textbox(
            MARGIN, r_top, Inches(1.3), Inches(0.3)
        )
        p = tf.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = row_label
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR
        run.font.name = "Inter"

        # Tool pills
        pill_left = MARGIN + Inches(1.4)
        for tool in tools:
            pill_w = Inches(max(0.8, len(tool) * 0.085 + 0.3))
            pill = add_badge(slide, pill_left, r_top + Inches(0.02),
                             tool, BG_CARD, TITLE_COLOR, 9)
            pill.line.color.rgb = CYAN
            pill.line.width = Pt(1)
            pill_left += pill_w + Inches(0.12)

    # --- Bottom insight card ---
    insight_top = row_top + Inches(1.15)
    insight_h = Inches(0.45)
    add_card(slide, MARGIN + Inches(1), insight_top,
             CONTENT_W - Inches(2), insight_h,
             fill_color=WHITE, border_color=GREEN, border_width=Pt(2))
    add_text_in_card(
        slide, MARGIN + Inches(1.2), insight_top + Inches(0.05),
        CONTENT_W - Inches(2.4), insight_h, [
            ("Um engenheiro + 6 agentes IA = pipeline completo de credit risk em < 7 horas", 14, TITLE_COLOR, True, False),
        ], alignment=PP_ALIGN.CENTER
    )


# =============================================================================
# SLIDE 17 — ROADMAP + ENCERRAMENTO
# =============================================================================
def create_slide_17(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, prs)
    add_footer(slide, 17)
    add_circuit_decoration(slide, prs)

    content_top = add_slide_title(slide, "Próximos Passos")

    # --- 3 Phase Timeline ---
    phases = [
        {
            "title": "Fase 1 — CURTO PRAZO (0-3 meses)",
            "badge": "DEPLOY", "badge_color": CYAN,
            "items": [
                "Deploy do champion em produção (batch scoring mensal)",
                "Integração com sistemas de aprovação Claro",
                "Alertas automáticos de drift (PSI > 0.10)",
            ],
            "border": CYAN,
        },
        {
            "title": "Fase 2 — MÉDIO PRAZO (3-6 meses)",
            "badge": "EVOLUÇÃO", "badge_color": GREEN,
            "items": [
                "Retreino trimestral com dados frescos",
                "A/B testing de cutoffs por segmento",
                "Score de crédito em tempo real (API REST)",
            ],
            "border": GREEN,
        },
        {
            "title": "Fase 3 — LONGO PRAZO (6-12 meses)",
            "badge": "EXPANSÃO", "badge_color": YELLOW, "badge_text": TITLE_COLOR,
            "items": [
                "Expansão para outros produtos Claro (fibra, TV)",
                "Modelo de lifetime value (LTV) + churn",
                "Federação de dados com parceiros",
            ],
            "border": YELLOW,
        },
    ]

    phase_w = CONTENT_W / 3 - Inches(0.12)
    phase_h = Inches(2.2)
    phase_top = content_top + Inches(0.15)

    for i, phase in enumerate(phases):
        p_left = MARGIN + i * (phase_w + Inches(0.18))

        add_card(slide, p_left, phase_top, phase_w, phase_h,
                 fill_color=WHITE, border_color=phase['border'], border_width=Pt(2))

        # Title
        tf = slide.shapes.add_textbox(
            p_left + Inches(0.15), phase_top + Inches(0.12),
            phase_w - Inches(0.3), Inches(0.35)
        )
        tf.text_frame.word_wrap = True
        p = tf.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = phase['title']
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR
        run.font.name = "Inter"

        # Badge
        badge_text_color = phase.get('badge_text', WHITE)
        add_badge(slide, p_left + Inches(0.15), phase_top + Inches(0.5),
                  phase['badge'], phase['badge_color'], badge_text_color, 10)

        # Items
        items_lines = [("• " + item, 10, BODY_COLOR, False, False)
                       for item in phase['items']]
        add_text_in_card(slide, p_left + Inches(0.15), phase_top + Inches(0.8),
                         phase_w - Inches(0.3), Inches(1.2), items_lines)

    # --- Impact card ---
    impact_top = phase_top + phase_h + Inches(0.25)
    impact_h = Inches(1.3)
    add_card(slide, MARGIN + Inches(0.5), impact_top,
             CONTENT_W - Inches(1), impact_h,
             fill_color=WHITE, border_color=CYAN, border_width=Pt(2))

    # 4 KPIs inline
    kpi_data = [
        ("3", "Cenários de cutoff", "ferramenta de decisão"),
        ("+R$ 1.7M líq.", "Cenário agressivo", "cutoff 300"),
        ("< R$ 170/mês", "Custo OCI", "infraestrutura"),
        ("3.9M", "Clientes scorados", "cobertura"),
    ]

    kpi_w = (CONTENT_W - Inches(1.4)) / 4
    for i, (val, label, sub) in enumerate(kpi_data):
        k_left = MARGIN + Inches(0.7) + i * kpi_w
        val_color = GREEN if "1.7M" in val else TITLE_COLOR

        add_text_in_card(slide, k_left, impact_top + Inches(0.1),
                         kpi_w, Inches(0.8), [
            (val, 24, val_color, True, False),
            (label, 10, TITLE_COLOR, True, False),
            (sub, 8, LABEL_COLOR, False, False),
        ], alignment=PP_ALIGN.CENTER)

    # Tagline
    tag_tf = slide.shapes.add_textbox(
        MARGIN + Inches(0.5), impact_top + Inches(0.85),
        CONTENT_W - Inches(1), Inches(0.2)
    )
    tag_tf.text_frame.word_wrap = True
    p = tag_tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "De dados brutos a decisões inteligentes — em 6 horas e 43 minutos."
    run.font.size = Pt(13)
    run.font.color.rgb = BODY_COLOR
    run.font.name = "Inter"

    # Disclaimer
    disc_tf = slide.shapes.add_textbox(
        MARGIN + Inches(0.5), impact_top + Inches(1.05),
        CONTENT_W - Inches(1), Inches(0.2)
    )
    disc_tf.text_frame.word_wrap = True
    p = disc_tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "*Economia líquida depende do cutoff e do LGD real. Modelo entrega a ferramenta; Claro define a política."
    run.font.size = Pt(9)
    run.font.color.rgb = LABEL_COLOR
    run.font.name = "Inter"

    # --- Obrigado ---
    obrigado_top = impact_top + impact_h + Inches(0.25)
    tf = slide.shapes.add_textbox(
        MARGIN, obrigado_top, CONTENT_W, Inches(0.5)
    )
    p = tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Obrigado."
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.name = "Inter"

    tf2 = slide.shapes.add_textbox(
        MARGIN, obrigado_top + Inches(0.55), CONTENT_W, Inches(0.3)
    )
    p2 = tf2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Equipe Hackathon PoD Academy — Claro + Oracle"
    run2.font.size = Pt(16)
    run2.font.color.rgb = BODY_COLOR
    run2.font.name = "Inter"


# =============================================================================
# MAIN
# =============================================================================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Generating Slide 14 — Squad OCI...")
    create_slide_14(prs)
    print("Generating Slide 15 — Mind Cloning...")
    create_slide_15(prs)
    print("Generating Slide 16 — Produtividade IA...")
    create_slide_16(prs)
    print("Generating Slide 17 — Roadmap + Encerramento...")
    create_slide_17(prs)

    output_path = "scripts/presentation/v2-batch-4-slides-14-17.pptx"
    prs.save(output_path)
    print(f"\nDone! Saved to: {output_path}")
    print("Open in PowerPoint to verify the design system.")


if __name__ == "__main__":
    main()
